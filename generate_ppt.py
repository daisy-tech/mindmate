from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(20, 39, 78)
BLUE = RGBColor(52, 111, 191)
LIGHT_BLUE = RGBColor(227, 240, 255)
GRAY = RGBColor(92, 104, 120)
LIGHT = RGBColor(248, 250, 252)
ACCENT = RGBColor(0, 168, 150)


def add_title(slide, title, subtitle=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT

    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = NAVY
    rect.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(8.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(9.5), Inches(0.4))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.runs[0]
        run2.font.size = Pt(12)
        run2.font.color.rgb = GRAY


def add_bullets(slide, bullets, left=0.7, top=1.4, width=5.8, height=5.4, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = NAVY
        if i == 0:
            p.bullet = True


def add_card(slide, x, y, w, h, title, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(224, 232, 240)
    shape.line.width = Pt(1)

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.3), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.bold = True
        run.font.size = Pt(17)
        run.font.color.rgb = BLUE

    tb2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(w - 0.3), Inches(h - 0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body
    p2.alignment = PP_ALIGN.LEFT
    for run in p2.runs:
        run.font.size = Pt(13)
        run.font.color.rgb = GRAY

# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
# background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = LIGHT

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.line.fill.background()

# top banner
banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(2.1))
banner.fill.solid()
banner.fill.fore_color.rgb = NAVY
banner.line.fill.background()

# title text
textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(8.8), Inches(1.1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = 'MindMate 伴学机器人'
p.alignment = PP_ALIGN.LEFT
for run in p.runs:
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(9.5), Inches(0.6))
tf2 = sub.text_frame
p2 = tf2.paragraphs[0]
p2.text = '把“AI 学习助手”升级为“持续陪伴的学习伙伴”'
p2.alignment = PP_ALIGN.LEFT
for run in p2.runs:
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(230, 240, 255)

# right side cards
add_card(slide, 9.0, 2.2, 3.2, 1.3, '核心定位', '教育 + 陪伴 + 记忆')
add_card(slide, 9.0, 3.95, 3.2, 1.3, '用户价值', '陪伴式学习，降低焦虑，提升效率')
add_card(slide, 9.0, 5.7, 3.2, 1.3, '愿景', '成为孩子成长路上的智能伴学伙伴')

# bottom features
add_card(slide, 0.8, 2.7, 3.0, 1.2, '听得懂', '自然对话，语音互动')
add_card(slide, 4.2, 2.7, 3.0, 1.2, '会学习', '根据年龄和水平做个性化教学')
add_card(slide, 7.6, 2.7, 3.0, 1.2, '有记忆', '记住孩子兴趣与学习进展')

# Slide 2: Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '问题是什么', '孩子学习遇到的真实痛点')
add_bullets(slide, [
    '传统辅导成本高，且难以做到 24 小时陪伴；家长难以持续参与学习过程。',
    '孩子注意力分散，学习内容碎片化，缺少“像朋友一样”的互动反馈。',
    '一对一教学难以规模化，孩子的学习节奏和情绪状态也很难被持续照顾。'
], left=0.8, top=1.6, width=6.0, height=4.8, font_size=19)

# add right visual
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.75), Inches(4.8), Inches(4.6))
shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.color.rgb = BLUE

box = slide.shapes.add_textbox(Inches(7.95), Inches(2.0), Inches(4.2), Inches(3.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '“孩子不只是要答案，更需要被理解、被鼓励、被陪伴。”'
for run in p.runs:
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = NAVY

# Slide 3: Market opportunity
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '为什么现在是好时机', '教育智能化与情感陪伴需求同时上升')
add_card(slide, 0.8, 1.7, 3.8, 1.8, '教育升级', 'AI 教学工具正在快速普及，家长愿意为更高效、更个性化的学习体验付费。')
add_card(slide, 4.9, 1.7, 3.8, 1.8, '陪伴需求', '儿童对“有情感连接”的智能产品接受度高，沉浸式互动能显著提升参与度。')
add_card(slide, 8.9, 1.7, 3.8, 1.8, '硬件入口', '智能音箱、桌面机器人、语音交互设备正形成新的家庭场景入口。')

# slide 4 solution
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '我们的产品是什么', '一款能“听、说、记、学”的伴学机器人')
add_bullets(slide, [
    '自然语音对话：支持日常问答、作业辅导、故事讲述。',
    '个性化学习：根据年龄、水平和兴趣制定学习节奏。',
    '记忆与长期陪伴：记住孩子偏好、进步和情绪状态，形成持续连接。',
    '家长视角：提供学习记录、成长反馈和亲子互动入口。'
], left=0.8, top=1.6, width=6.1, height=4.8, font_size=18)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.75), Inches(4.8), Inches(4.6))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(245, 249, 255)
shape.line.color.rgb = BLUE
box = slide.shapes.add_textbox(Inches(7.8), Inches(2.1), Inches(4.1), Inches(3.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '核心能力：\n1）自然交互\n2）长期记忆\n3）学习个性化\n4）情感陪伴'
for run in p.runs:
    run.font.size = Pt(22)
    run.font.color.rgb = NAVY

# Slide 5: use cases
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '它能解决哪些场景', '从课后学习到睡前陪伴')
add_card(slide, 0.8, 1.8, 3.8, 2.1, '课后作业', '孩子遇到问题时，机器人即时解答并引导思考。')
add_card(slide, 4.9, 1.8, 3.8, 2.1, '睡前故事', '讲故事、讲知识、安抚情绪，增强亲子连接。')
add_card(slide, 8.9, 1.8, 3.8, 2.1, '日常陪伴', '陪孩子聊天、复习、打卡、做小游戏，建立稳定习惯。')
add_card(slide, 0.8, 4.4, 11.7, 1.6, '一句话价值', '不是“替孩子学习”，而是“成为孩子学习路上的智能朋友”。')

# Slide 6: business model
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '商业模式', '从硬件、内容到服务的复合增长')
add_card(slide, 0.8, 1.8, 3.8, 2.1, '硬件销售', '机器人终端作为家庭入口，形成品牌与场景认知。')
add_card(slide, 4.9, 1.8, 3.8, 2.1, '订阅服务', '学习内容、个性化陪伴、成长分析等按月付费。')
add_card(slide, 8.9, 1.8, 3.8, 2.1, 'B2B 扩展', '面向学校、培训机构和亲子平台输出定制化解决方案。')
add_card(slide, 0.8, 4.5, 11.7, 1.6, '收入逻辑', '先做用户认知，再靠高频使用和长期续费形成稳定收入。')

# Slide 7: differentiation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '为什么我们有机会赢', '差异化竞争优势')
add_bullets(slide, [
    '更贴近儿童使用习惯：语音、情绪、游戏化更容易被接受。',
    '具备长期记忆能力：不是一次性问答，而是持续成长的“陪伴关系”。',
    '有明确的家庭场景入口：能与家长、孩子、学校形成闭环。',
    '从“工具型AI”升级为“有陪伴感的智能伙伴”。'
], left=0.8, top=1.6, width=6.2, height=4.8, font_size=18)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.75), Inches(4.8), Inches(4.6))
shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.color.rgb = BLUE
box = slide.shapes.add_textbox(Inches(7.8), Inches(2.1), Inches(4.1), Inches(3.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '差异点：\n- 不是纯聊天工具\n- 不是简单教育软件\n- 而是“有记忆、有陪伴、有成长”的机器人产品'
for run in p.runs:
    run.font.size = Pt(20)
    run.font.color.rgb = NAVY

# Slide 8: roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '产品路线图', '从原型验证到规模化落地')
add_card(slide, 0.8, 1.8, 3.8, 2.1, '第一阶段', '搭建原型，完成语音对话、学习问答和基础记忆功能。')
add_card(slide, 4.9, 1.8, 3.8, 2.1, '第二阶段', '完成亲子场景测试，优化内容与交互体验。')
add_card(slide, 8.9, 1.8, 3.8, 2.1, '第三阶段', '发力订阅服务、教育合作和家庭场景扩展。')
add_card(slide, 0.8, 4.5, 11.7, 1.6, '目标', '用真实用户反馈不断降低成本、提高留存和付费转化。')

# Slide 9: validation KPI
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '如何验证', '用真实使用数据证明产品价值')
add_bullets(slide, [
    '用户留存：首周、首月活跃和复访率。',
    '学习效果：作业完成率、知识点掌握反馈、学习习惯建立。',
    '付费意愿：订阅转化、复购与家庭渗透率。',
    '体验反馈：陪伴感评分、家长满意度、孩子使用时长。'
], left=0.8, top=1.6, width=6.2, height=4.8, font_size=18)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.75), Inches(4.8), Inches(4.6))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(247, 251, 255)
shape.line.color.rgb = BLUE
box = slide.shapes.add_textbox(Inches(7.8), Inches(2.1), Inches(4.1), Inches(3.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '关键指标：\n- 日活 / 周活\n- 复购率\n- 家长评分\n- 学习完成率'
for run in p.runs:
    run.font.size = Pt(22)
    run.font.color.rgb = NAVY

# Slide 10: closing/ask
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '结语', '让 AI 走进孩子的成长，成为可信赖的学习伙伴')
add_bullets(slide, [
    'MindMate 的核心不是“更像一个工具”，而是“更像一个会陪伴的伙伴”。',
    '我们相信，未来的教育产品将是“有记忆、有情感、有反馈”的智能陪伴。',
    '如果你认可这个方向，我们希望与你一起把它做成真正落地的产品。'
], left=0.8, top=1.8, width=8.7, height=3.8, font_size=20)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.9), Inches(1.95), Inches(2.7), Inches(2.6))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
box = slide.shapes.add_textbox(Inches(10.15), Inches(2.35), Inches(2.2), Inches(1.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '合作/投资/试点'
p.alignment = PP_ALIGN.CENTER
for run in p.runs:
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

prs.save('mindmate_路演PPT.pptx')
print('saved to mindmate_路演PPT.pptx')
